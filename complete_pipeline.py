from pptx import Presentation
import pke
import spacy
from summarizer import Summarizer as Sum
import itertools
import pandas as pd
from tqdm import tqdm
import time
import csv
import json
import requests

def process_data(folder_name,course_name,course_intro,courseId):
    fo = open(course_name+"content.txt", "w",encoding='utf8', errors='ignore')

    for filein in os.listdir(folder_name):
        prs = Presentation(os.path.join(folder_name,filein))
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        for slide in prs.slides:
            list_of_elements = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    line = ''
                    for run in paragraph.runs:
                        line += run.text	
                    list_of_elements.append(line)
                
            for elements in list_of_elements:
                print(elements)
                fo.write(elements +'\n')
            fo.write('\n')


    t1 = time.perf_counter()
    f = open(course_name+"content.txt",'r',encoding='utf8', errors='ignore')
    text = f.read()
    text = text.replace('\n','.')
    nlp = spacy.load("en_core_web_lg")




    def keyphrase_extractor(text,num=10):
        # initialize keyphrase extraction model, here TopicRank
        extractor = pke.unsupervised.TopicRank()
        # load the content of the document, here document is expected to be in raw
        # format (i.e. a simple text file) and preprocessing is carried out using spacy
        extractor.load_document(input=text, language='en')
        # keyphrase candidate selection, in the case of TopicRank: sequences of nouns
        # and adjectives (i.e. `(Noun|Adj)*`)
        extractor.candidate_selection()

        # candidate weighting, in the case of TopicRank: using a random walk algorithm
        extractor.candidate_weighting(threshold=0.74, method='average')

        # N-best selection, keyphrases contains the 10 highest scored candidates as
        # (keyphrase, score) tuples
        keyphrases = extractor.get_n_best(n=num)
        return keyphrases

    def search_doc_for_keyphrases(text,keys,nlp):
        doc1 = nlp(text)
        sents = list(doc1.sents)
        key_sentence_dict = dict()
        for keyphrase in tqdm(keys):
            popped = 0

            key_sentence_dict[keyphrase] = []
            for ind,sent in enumerate(sents):
                sent = str(sent)
                sent= sent.lower()
                if keyphrase in sent:
                    key_sentence_dict[keyphrase].append(sent)
            key_sentence_dict[keyphrase] = " ".join(key_sentence_dict[keyphrase])
        return key_sentence_dict,sents

    def summarizer(key_sentence_dict):
        summary = Sum()
        summarized = dict()
        for keys,text in tqdm(key_sentence_dict.items()):
            summarized[keys] = summary(text)
        return summarized

    def similarity_generator(summarized,sents,nlp):
        similar = dict()
        for key,summary in tqdm(summarized.items()):
            similar[key] = []
            sum_doc = nlp(summary)
            x=0
            temp_set = set()
            for sent in sents:
                sent = str(sent)
                sent = sent.replace('.','')
                if sent in temp_set:
                    continue
                temp_set.add(sent)
                temp_doc = nlp(sent)
                sim = sum_doc.similarity(temp_doc)
                similar[key].append((sent,sim,x))
                x+=1
        return similar

    def semantic_summary_extractor(similar,num_sentences = 8):
        semantic_summary = dict()
        for key,data in tqdm(similar.items()):
            sorted_data_similarity = sorted(data,key=lambda x: x[1])
            sorted_data_similarity=sorted_data_similarity[::-1]
            selected_sorted = sorted_data_similarity[:num_sentences]
            final_sorted = sorted(selected_sorted,key=lambda x: x[2])
            final_sorted = [x[0] for x in final_sorted]
            semantic_summary[key] = ". ".join(final_sorted)
        return semantic_summary

    keys = []
    keyphrases = keyphrase_extractor('content.txt',num = 50)
    for key, weight in keyphrases:
        keys.append(key)
        print(key)

    ks,sents = search_doc_for_keyphrases(text,keys,nlp)
    print(ks,sents)
    summarized = summarizer(ks)
    print(summarized)
    similar = similarity_generator(summarized,sents,nlp)
    semantic_sum = semantic_summary_extractor(similar,20)
    df = pd.DataFrame.from_dict(semantic_sum, orient="index")
    df.to_csv(course_name+"content.csv")
    t2 = time.perf_counter()
    print("Time taken: {} seconds".format(t2-t1))

    # # for topic in list(topic_summary_dict.keys()):
    # #     print(topic_summary_dict[topic])
    topic_summary_dict = {}

    topic_summary_dict = semantic_sum

    topics = list(topic_summary_dict.keys())

    topic_topic_weights = {}

    for key,value in tqdm(topic_summary_dict.items()):
        topic_topic_weights[key] = []
        curr_doc = nlp(value)
        for topic in topics:
            temp_doc = nlp(topic_summary_dict[topic])
            topic_topic_weights[key].append((topic,curr_doc.similarity(temp_doc)))
        topic_topic_weights[key] = sorted(topic_topic_weights[key], key= lambda x:x[1],reverse= True)
        topic_topic_weights[key] = topic_topic_weights[key][:10]
    print(topic_topic_weights)

    topics = {}
    for topic,summary in topic_summary_dict.items():
        topics[topic] ={}
        topics[topic]["summary"]=summary
        topics[topic]["links"] = {}
        for temp_topic,weight in topic_topic_weights[topic]:
            topics[topic]["links"][course_name+"/"+temp_topic] = str(weight)
        topics[topic]["questions"] = ["q1","q2"]

    course = {
        "courseId":courseId,
        "data":{
            "description":course_intro,
            "name":course_name,
        
        "topics":{
            "topicName":{
                "links":{
                    "course/topic":"weight"
                },
                "questions":[""],
                "summary":""
            }
        }
            }
    }

    course["data"]["topics"] = topics


    a = json.dumps(course)


    url = "https://us-central1-edugraph-78964.cloudfunctions.net/app/courses/upload/"

    payload = a
    headers = {
    'Content-Type': 'application/javascript'
    }

    response = requests.request("POST", url, headers=headers, data = payload)

    print(response.text.encode('utf8'))